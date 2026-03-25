from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
OUTPUT_DIR = ROOT / "docs" / "presentations"
OUTPUT_PATH = OUTPUT_DIR / "mango-internal-architecture-report.pptx"


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

MANGO = RGBColor(245, 166, 35)
MANGO_DARK = RGBColor(181, 107, 23)
MANGO_LIGHT = RGBColor(255, 239, 213)
INK = RGBColor(32, 25, 20)
INK_SOFT = RGBColor(92, 79, 67)
BG = RGBColor(252, 248, 241)
PANEL = RGBColor(255, 255, 255)
LINE = RGBColor(222, 209, 190)
SUCCESS = RGBColor(56, 142, 96)
WARN = RGBColor(205, 124, 24)
RISK = RGBColor(176, 55, 55)
INFO = RGBColor(53, 104, 173)

TITLE_FONT = "Aptos Display"
BODY_FONT = "Aptos"
MONO_FONT = "Consolas"


def set_slide_background(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_top_bar(slide, title: str, subtitle: str | None = None):
    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.5)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = MANGO
    band.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.68), Inches(8.8), Inches(0.6))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = TITLE_FONT
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = INK

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.62), Inches(1.22), Inches(8.8), Inches(0.35))
        stf = sub.text_frame
        stf.clear()
        p = stf.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        r.font.name = BODY_FONT
        r.font.size = Pt(11.5)
        r.font.color.rgb = INK_SOFT


def add_footer(slide, text: str):
    footer = slide.shapes.add_textbox(Inches(0.6), Inches(7.02), Inches(12.0), Inches(0.22))
    tf = footer.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = text
    r.font.name = BODY_FONT
    r.font.size = Pt(9)
    r.font.color.rgb = INK_SOFT


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size=18,
    bold=False,
    color=INK,
    font_name=BODY_FONT,
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font_name
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def add_bullets(
    slide,
    left,
    top,
    width,
    height,
    bullets,
    font_size=16,
    color=INK,
    level0_bullet="•",
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    for index, bullet in enumerate(bullets):
        if index == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = f"{level0_bullet} {bullet}"
        r.font.name = BODY_FONT
        r.font.size = Pt(font_size)
        r.font.color.rgb = color
    return box


def add_panel(slide, left, top, width, height, title=None, fill=PANEL, line=LINE, radius=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    panel = slide.shapes.add_shape(shape_type, left, top, width, height)
    panel.fill.solid()
    panel.fill.fore_color.rgb = fill
    panel.line.color.rgb = line
    panel.line.width = Pt(1)
    if title:
        add_textbox(
            slide,
            left + Inches(0.18),
            top + Inches(0.1),
            width - Inches(0.36),
            Inches(0.3),
            title,
            font_size=13,
            bold=True,
            color=INK_SOFT,
        )
    return panel


def add_stat_card(slide, left, top, width, height, number, label, accent=MANGO):
    add_panel(slide, left, top, width, height, fill=PANEL, line=LINE)
    add_textbox(
        slide,
        left + Inches(0.18),
        top + Inches(0.15),
        width - Inches(0.36),
        Inches(0.42),
        number,
        font_size=24,
        bold=True,
        color=accent,
    )
    add_textbox(
        slide,
        left + Inches(0.18),
        top + Inches(0.62),
        width - Inches(0.36),
        Inches(0.45),
        label,
        font_size=11.5,
        color=INK_SOFT,
    )


def add_box_with_title(slide, left, top, width, height, title, lines, fill=PANEL, accent=MANGO):
    add_panel(slide, left, top, width, height, fill=fill, line=accent)
    title_band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.38)
    )
    title_band.fill.solid()
    title_band.fill.fore_color.rgb = accent
    title_band.line.fill.background()
    title_tf = title_band.text_frame
    title_tf.clear()
    title_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = title_tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.name = BODY_FONT
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = PANEL
    add_bullets(
        slide,
        left + Inches(0.15),
        top + Inches(0.48),
        width - Inches(0.3),
        height - Inches(0.6),
        lines,
        font_size=11.5,
    )


def connect_right(slide, from_shape, to_shape, color=INK_SOFT):
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        from_shape.left + from_shape.width,
        from_shape.top + from_shape.height / 2,
        to_shape.left,
        to_shape.top + to_shape.height / 2,
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(1.5)
    connector.line.end_arrowhead = True


def connect_down(slide, from_shape, to_shape, color=INK_SOFT):
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        from_shape.left + from_shape.width / 2,
        from_shape.top + from_shape.height,
        to_shape.left + to_shape.width / 2,
        to_shape.top,
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(1.5)
    connector.line.end_arrowhead = True


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    title_layout = prs.slide_layouts[6]

    # Slide 1: Cover
    slide = prs.slides.add_slide(title_layout)
    set_slide_background(slide)

    hero = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.62), Inches(0.6), Inches(12.05), Inches(5.7)
    )
    hero.fill.solid()
    hero.fill.fore_color.rgb = PANEL
    hero.line.color.rgb = LINE
    hero.line.width = Pt(1)

    mango_band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.62), Inches(0.6), Inches(12.05), Inches(0.56)
    )
    mango_band.fill.solid()
    mango_band.fill.fore_color.rgb = MANGO
    mango_band.line.fill.background()

    add_textbox(slide, Inches(0.95), Inches(1.45), Inches(8.8), Inches(0.65), "Mango 架构与研发现状汇报", 28, True)
    add_textbox(
        slide,
        Inches(0.98),
        Inches(2.15),
        Inches(8.7),
        Inches(0.45),
        "团队内部研发汇报 · 基于 codex-mango-v1 当前实现与文档",
        14,
        False,
        INK_SOFT,
    )
    add_textbox(
        slide,
        Inches(0.98),
        Inches(3.0),
        Inches(5.7),
        Inches(1.2),
        "核心结论：\nMango 已完成桌面端 Agent 工作台骨架，主闭环已打通，当前重点从“能跑”进入“真实可用”阶段。",
        18,
        False,
        INK,
    )

    add_stat_card(slide, Inches(8.55), Inches(1.45), Inches(1.45), Inches(1.25), "4", "主流程阶段")
    add_stat_card(slide, Inches(10.15), Inches(1.45), Inches(1.45), Inches(1.25), "5", "核心共享包")
    add_stat_card(slide, Inches(8.55), Inches(2.95), Inches(1.45), Inches(1.25), "2", "当前可用 Adapter")
    add_stat_card(slide, Inches(10.15), Inches(2.95), Inches(1.45), Inches(1.25), "1", "正式主存储")

    add_bullets(
        slide,
        Inches(8.2),
        Inches(4.65),
        Inches(3.4),
        Inches(1.1),
        [
            "桌面端本地执行优先",
            "Plan -> Approve -> Go -> Review",
            "SQLite 主存储已落地",
        ],
        font_size=12,
    )
    add_footer(slide, "Mango · Internal Engineering Report · 2026-03-25")

    # Slide 2: Executive summary
    slide = prs.slides.add_slide(title_layout)
    set_slide_background(slide)
    add_top_bar(slide, "1. 项目定位与当前判断", "先用一页把“我们现在到底做到哪”讲清楚")

    add_box_with_title(
        slide,
        Inches(0.7),
        Inches(1.6),
        Inches(3.8),
        Inches(2.2),
        "产品定位",
        [
            "面向开发者的桌面端 Agent 工作台",
            "不是聊天产品，而是执行产品",
            "强调可控、可信、可追踪、可复盘",
        ],
    )
    add_box_with_title(
        slide,
        Inches(4.78),
        Inches(1.6),
        Inches(3.8),
        Inches(2.2),
        "当前阶段",
        [
            "M1 可运行骨架已完成",
            "M2 真实执行闭环正在形成",
            "核心主链路已可本地验证",
        ],
    )
    add_box_with_title(
        slide,
        Inches(8.86),
        Inches(1.6),
        Inches(3.8),
        Inches(2.2),
        "本轮判断",
        [
            "架构边界清楚",
            "真实 CLI 接入已开始",
            "执行深度和控制精度仍需增强",
        ],
    )

    add_panel(slide, Inches(0.7), Inches(4.15), Inches(11.96), Inches(2.1), "一句话结论")
    add_textbox(
        slide,
        Inches(0.95),
        Inches(4.6),
        Inches(11.4),
        Inches(1.15),
        "Mango 已不再是纯规划项目，而是一套具备桌面壳、领域模型、IPC 契约、适配器层、SQLite 主存储和基础质量门禁的可演进架构骨架；但它仍需继续补强真实执行事件、细粒度权限控制、恢复重试和更强的回顾能力，才能从“骨架可跑”进入“团队可用 / Beta 可试用”。",
        17,
    )
    add_footer(slide, "Source: README / PRD / 路线图 / 详细架构设计")

    # Slide 3: System architecture
    slide = prs.slides.add_slide(title_layout)
    set_slide_background(slide)
    add_top_bar(slide, "2. 当前整体架构", "仓库分层 + 运行时分层同时存在，核心是边界清楚")

    a = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(2.25), Inches(2.1), Inches(1.0)
    )
    a.fill.solid()
    a.fill.fore_color.rgb = MANGO_LIGHT
    a.line.color.rgb = MANGO_DARK
    add_textbox(slide, Inches(1.1), Inches(2.45), Inches(1.7), Inches(0.4), "apps/desktop", 18, True, INK, TITLE_FONT, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.95), Inches(3.35), Inches(2.0), Inches(0.45), "当前唯一完整运行入口", 11.5, False, INK_SOFT, BODY_FONT, PP_ALIGN.CENTER)

    b = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(3.45), Inches(2.25), Inches(2.15), Inches(1.0)
    )
    b.fill.solid()
    b.fill.fore_color.rgb = PANEL
    b.line.color.rgb = LINE
    add_textbox(slide, Inches(3.7), Inches(2.42), Inches(1.65), Inches(0.48), "packages/core", 18, True, INK, TITLE_FONT, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(3.53), Inches(3.35), Inches(1.98), Inches(0.45), "状态机 / 权限 / 事件 / 工作区", 11, False, INK_SOFT, BODY_FONT, PP_ALIGN.CENTER)

    c = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.1), Inches(2.25), Inches(2.15), Inches(1.0)
    )
    c.fill.solid()
    c.fill.fore_color.rgb = PANEL
    c.line.color.rgb = LINE
    add_textbox(slide, Inches(6.28), Inches(2.42), Inches(1.8), Inches(0.48), "packages/adapters", 18, True, INK, TITLE_FONT, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(6.18), Inches(3.35), Inches(1.98), Inches(0.45), "Claude CLI / Mock", 11.5, False, INK_SOFT, BODY_FONT, PP_ALIGN.CENTER)

    d = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.75), Inches(2.25), Inches(1.95), Inches(1.0)
    )
    d.fill.solid()
    d.fill.fore_color.rgb = PANEL
    d.line.color.rgb = LINE
    add_textbox(slide, Inches(8.95), Inches(2.42), Inches(1.55), Inches(0.48), "contracts", 18, True, INK, TITLE_FONT, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(8.83), Inches(3.35), Inches(1.75), Inches(0.45), "IPC DTO / 错误模型", 11.5, False, INK_SOFT, BODY_FONT, PP_ALIGN.CENTER)

    e = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(11.0), Inches(2.25), Inches(1.4), Inches(1.0)
    )
    e.fill.solid()
    e.fill.fore_color.rgb = PANEL
    e.line.color.rgb = LINE
    add_textbox(slide, Inches(11.14), Inches(2.42), Inches(1.1), Inches(0.48), "SQLite", 18, True, INK, TITLE_FONT, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(10.98), Inches(3.35), Inches(1.25), Inches(0.45), "正式主存储", 11.5, False, INK_SOFT, BODY_FONT, PP_ALIGN.CENTER)

    connect_right(slide, a, b)
    connect_right(slide, b, c)
    connect_right(slide, c, d)
    connect_right(slide, d, e)

    add_panel(slide, Inches(0.9), Inches(4.45), Inches(11.5), Inches(1.55), "架构判断")
    add_bullets(
        slide,
        Inches(1.1),
        Inches(4.82),
        Inches(11.1),
        Inches(0.95),
        [
            "当前最重要的不是模块数量，而是 Desktop / Core / Adapters / Contracts 的边界已经成立",
            "未来 Web、API、Worker 已预留位置，但不会抢走当前桌面本地主链路的优先级",
            "这套结构足以支撑从 mock 走向真实 CLI，再走向更完整的恢复、回顾和多 adapter 能力",
        ],
        font_size=13.5,
    )
    add_footer(slide, "Source: technical-architecture.md / module-contracts.md")

    # Slide 4: Desktop runtime
    slide = prs.slides.add_slide(title_layout)
    set_slide_background(slide)
    add_top_bar(slide, "3. 桌面端运行时架构", "当前唯一完整运行路径：Renderer -> Preload -> Main -> Adapter / Store")

    renderer = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(2.1), Inches(2.2), Inches(1.15)
    )
    renderer.fill.solid()
    renderer.fill.fore_color.rgb = PANEL
    renderer.line.color.rgb = INFO
    add_textbox(slide, Inches(1.2), Inches(2.4), Inches(1.55), Inches(0.32), "Renderer", 20, True, INFO, TITLE_FONT, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.98), Inches(3.35), Inches(1.95), Inches(0.6), "React UI\n任务输入 / 计划展示 / 时间线 / 回顾", 11.5, False, INK_SOFT, BODY_FONT, PP_ALIGN.CENTER)

    preload = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(3.45), Inches(2.1), Inches(2.2), Inches(1.15)
    )
    preload.fill.solid()
    preload.fill.fore_color.rgb = PANEL
    preload.line.color.rgb = MANGO_DARK
    add_textbox(slide, Inches(3.88), Inches(2.4), Inches(1.35), Inches(0.32), "Preload", 20, True, MANGO_DARK, TITLE_FONT, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(3.58), Inches(3.35), Inches(1.95), Inches(0.6), "DesktopApi\n只做安全桥和类型化暴露", 11.5, False, INK_SOFT, BODY_FONT, PP_ALIGN.CENTER)

    main = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.05), Inches(2.1), Inches(2.2), Inches(1.15)
    )
    main.fill.solid()
    main.fill.fore_color.rgb = MANGO_LIGHT
    main.line.color.rgb = MANGO_DARK
    add_textbox(slide, Inches(6.45), Inches(2.4), Inches(1.45), Inches(0.32), "Main", 20, True, MANGO_DARK, TITLE_FONT, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(6.16), Inches(3.35), Inches(1.98), Inches(0.6), "DesktopController\n任务编排 / IPC / 持久化 / 系统交互", 11.5, False, INK_SOFT, BODY_FONT, PP_ALIGN.CENTER)

    adapter = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.75), Inches(1.6), Inches(1.7), Inches(0.95)
    )
    adapter.fill.solid()
    adapter.fill.fore_color.rgb = PANEL
    adapter.line.color.rgb = SUCCESS
    add_textbox(slide, Inches(9.0), Inches(1.92), Inches(1.2), Inches(0.25), "Adapter", 18, True, SUCCESS, TITLE_FONT, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(8.88), Inches(2.55), Inches(1.45), Inches(0.5), "Claude CLI\n或 Mock", 11, False, INK_SOFT, BODY_FONT, PP_ALIGN.CENTER)

    store = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.75), Inches(3.05), Inches(1.7), Inches(0.95)
    )
    store.fill.solid()
    store.fill.fore_color.rgb = PANEL
    store.line.color.rgb = INFO
    add_textbox(slide, Inches(9.05), Inches(3.35), Inches(1.1), Inches(0.25), "Store", 18, True, INFO, TITLE_FONT, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(8.85), Inches(4.0), Inches(1.48), Inches(0.45), "SQLite 主存储\nJSON 导入兜底", 11, False, INK_SOFT, BODY_FONT, PP_ALIGN.CENTER)

    osbox = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(10.85), Inches(2.3), Inches(1.7), Inches(0.95)
    )
    osbox.fill.solid()
    osbox.fill.fore_color.rgb = PANEL
    osbox.line.color.rgb = WARN
    add_textbox(slide, Inches(11.15), Inches(2.6), Inches(1.1), Inches(0.25), "System", 18, True, WARN, TITLE_FONT, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(10.98), Inches(3.25), Inches(1.4), Inches(0.5), "Shell / 文件 / Git\n高风险能力入口", 11, False, INK_SOFT, BODY_FONT, PP_ALIGN.CENTER)

    connect_right(slide, renderer, preload)
    connect_right(slide, preload, main)
    connect_right(slide, main, osbox)
    connect_right(slide, main, adapter)
    connect_right(slide, main, store)

    add_panel(slide, Inches(0.85), Inches(5.2), Inches(11.7), Inches(1.1), "关键约束")
    add_bullets(
        slide,
        Inches(1.05),
        Inches(5.52),
        Inches(11.2),
        Inches(0.52),
        [
            "Renderer 不能直接碰 Node/Electron；业务编排不能下沉到 Preload；高风险能力必须由 Main 统一收口",
        ],
        font_size=14,
    )
    add_footer(slide, "Source: desktop runtime code / DesktopController / DesktopApi")

    # Slide 5: Domain model
    slide = prs.slides.add_slide(title_layout)
    set_slide_background(slide)
    add_top_bar(slide, "4. 核心领域模型与状态机", "UI、持久化和执行都围绕同一组对象组织")

    add_box_with_title(
        slide,
        Inches(0.75),
        Inches(1.55),
        Inches(3.15),
        Inches(2.45),
        "核心对象",
        [
            "TaskSession",
            "TaskPlan",
            "ExecutionEvent",
            "TaskReview",
            "WorkspaceContext",
            "PermissionPolicy",
        ],
        accent=INFO,
    )

    state_boxes = []
    labels = [("draft", INFO), ("planned", MANGO_DARK), ("approved", WARN), ("running", SUCCESS), ("succeeded / failed", RISK)]
    start_left = Inches(4.35)
    for idx, (label, color) in enumerate(labels):
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            start_left + Inches(1.6 * idx),
            Inches(2.05),
            Inches(1.3),
            Inches(0.72),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = PANEL
        shape.line.color.rgb = color
        add_textbox(
            slide,
            shape.left + Inches(0.05),
            shape.top + Inches(0.16),
            Inches(1.2),
            Inches(0.25),
            label,
            13,
            True,
            color,
            BODY_FONT,
            PP_ALIGN.CENTER,
        )
        state_boxes.append(shape)

    for idx in range(len(state_boxes) - 1):
        connect_right(slide, state_boxes[idx], state_boxes[idx + 1])

    add_textbox(slide, Inches(4.35), Inches(3.15), Inches(7.2), Inches(0.35), "TaskSession 固定流转：不允许跳过关键阶段", 14, True, INK_SOFT)

    add_box_with_title(
        slide,
        Inches(4.35),
        Inches(4.05),
        Inches(3.85),
        Inches(1.7),
        "ExecutionEvent",
        [
            "terminal.output",
            "file.change",
            "summary.ready",
            "tool.call",
        ],
        accent=SUCCESS,
    )
    add_box_with_title(
        slide,
        Inches(8.45),
        Inches(4.05),
        Inches(3.65),
        Inches(1.7),
        "Permission Capability",
        [
            "shell",
            "filesystem",
            "network",
            "browser",
        ],
        accent=RISK,
    )
    add_footer(slide, "Source: packages/core/src/types.ts / taskSession.ts")

    # Slide 6: Main flow
    slide = prs.slides.add_slide(title_layout)
    set_slide_background(slide)
    add_top_bar(slide, "5. 主执行链路", "当前主链路已经贯通，但执行事件和控制粒度仍是下一阶段重点")

    steps = [
        ("1. 输入目标", "用户描述任务目标"),
        ("2. 生成计划", "Adapter 返回结构化 TaskPlan"),
        ("3. 权限评估", "Main 计算风险与审批信息"),
        ("4. 批准执行", "任务进入 running"),
        ("5. 回写事件", "事件、状态和 review 写回存储"),
    ]

    shapes = []
    for idx, (title, desc) in enumerate(steps):
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(0.75 + idx * 2.45),
            Inches(2.05),
            Inches(2.0),
            Inches(1.55),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = PANEL
        shape.line.color.rgb = MANGO if idx in (1, 3) else LINE
        add_textbox(slide, shape.left + Inches(0.14), shape.top + Inches(0.2), Inches(1.72), Inches(0.3), title, 14, True)
        add_textbox(slide, shape.left + Inches(0.14), shape.top + Inches(0.62), Inches(1.72), Inches(0.62), desc, 11.5, False, INK_SOFT)
        shapes.append(shape)

    for idx in range(len(shapes) - 1):
        connect_right(slide, shapes[idx], shapes[idx + 1], MANGO_DARK)

    add_panel(slide, Inches(0.75), Inches(4.35), Inches(12.0), Inches(1.55), "当前链路中的 4 个刻意简化")
    add_bullets(
        slide,
        Inches(0.98),
        Inches(4.72),
        Inches(11.55),
        Inches(0.85),
        [
            "批准仍然是任务级整体批准，不是逐 capability 授权",
            "执行事件一次性返回，不是流式推送",
            "真实 CLI 尚未稳定产出 file.change / tool.call",
            "取消、恢复、重试还没有完整产品化",
        ],
        font_size=13,
    )
    add_footer(slide, "Source: DesktopController / ClaudeCodeCliAdapter / detailed-architecture.md")

    # Slide 7: Implemented capabilities
    slide = prs.slides.add_slide(title_layout)
    set_slide_background(slide)
    add_top_bar(slide, "6. 当前已落地能力", "这页回答：哪些东西已经不是概念，而是代码里真实存在")

    add_box_with_title(
        slide,
        Inches(0.75),
        Inches(1.55),
        Inches(3.9),
        Inches(4.75),
        "已经落地",
        [
            "Electron + React + TypeScript 桌面壳",
            "Main / Preload / Renderer 分层",
            "TaskSession 状态机与事件模型",
            "DesktopApi 与 IPC 契约层",
            "Claude CLI + Mock 双 adapter",
            "SQLite 主存储与 JSON 自动导入",
            "lint / test / typecheck / build / smoke",
        ],
        accent=SUCCESS,
    )

    add_box_with_title(
        slide,
        Inches(4.9),
        Inches(1.55),
        Inches(3.75),
        Inches(4.75),
        "已经证明可行",
        [
            "可以本地启动桌面工作台",
            "可以基于真实 CLI 生成结构化计划",
            "可以审批后执行一轮端到端任务",
            "可以把任务、事件、历史写入正式存储",
            "可以通过统一状态快照驱动 UI",
        ],
        accent=INFO,
    )

    add_box_with_title(
        slide,
        Inches(8.9),
        Inches(1.55),
        Inches(3.45),
        Inches(4.75),
        "当前价值",
        [
            "已具备继续迭代的稳定骨架",
            "已从纯 mock 进入真实执行验证",
            "后续重点不再是“从 0 建壳”，而是补强能力深度和可用性",
        ],
        accent=MANGO_DARK,
    )
    add_footer(slide, "Source: README / package.json / current repository structure")

    # Slide 8: Risks
    slide = prs.slides.add_slide(title_layout)
    set_slide_background(slide)
    add_top_bar(slide, "7. 主要缺口与风险", "这页回答：为什么现在还不能把它当成成熟 Beta 系统")

    add_box_with_title(
        slide,
        Inches(0.75),
        Inches(1.55),
        Inches(3.8),
        Inches(2.05),
        "执行可见性不足",
        [
            "真实 CLI 仍缺少流式事件",
            "文件变化与工具调用证据不足",
            "Review 面板证据链偏弱",
        ],
        accent=RISK,
    )
    add_box_with_title(
        slide,
        Inches(4.8),
        Inches(1.55),
        Inches(3.8),
        Inches(2.05),
        "控制粒度不足",
        [
            "权限还是任务级整体批准",
            "缺少取消、恢复、失败重试",
            "工作区级授权范围未完成",
        ],
        accent=WARN,
    )
    add_box_with_title(
        slide,
        Inches(8.85),
        Inches(1.55),
        Inches(3.45),
        Inches(2.05),
        "产品完成度不足",
        [
            "工作区管理偏基础",
            "多 adapter 策略还较初级",
            "安装、更新、反馈闭环仍需打磨",
        ],
        accent=INFO,
    )

    add_panel(slide, Inches(0.75), Inches(4.05), Inches(11.55), Inches(1.85), "研发含义")
    add_bullets(
        slide,
        Inches(0.98),
        Inches(4.4),
        Inches(11.1),
        Inches(1.1),
        [
            "下一阶段重点不应再扩散到新平台，而是把桌面本地主链路做深、做稳、做得可诊断",
            "如果过早转向 Web / API / Worker，会放大当前执行链路尚未稳定的风险",
            "当前更适合内部持续验证和定向试用，不适合被表述成“功能已完整”",
        ],
        font_size=13.2,
    )
    add_footer(slide, "Source: v1-backlog.md / roadmap-and-milestones.md / detailed-architecture.md")

    # Slide 9: Roadmap
    slide = prs.slides.add_slide(title_layout)
    set_slide_background(slide)
    add_top_bar(slide, "8. 下一阶段建议路线", "按“先补真实可用，再谈平台扩展”的顺序推进")

    roadmap_items = [
        ("P1", "补强真实 CLI 事件能力", SUCCESS),
        ("P2", "把权限从展示升级为真正控制点", RISK),
        ("P3", "补取消 / 恢复 / 失败重试", WARN),
        ("P4", "增强工作区管理与 Review 证据", INFO),
        ("P5", "再扩到 Web / API / Worker", MANGO_DARK),
    ]

    for idx, (tag, label, color) in enumerate(roadmap_items):
        left = Inches(0.9 + idx * 2.45)
        circle = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL, left + Inches(0.55), Inches(1.9), Inches(0.78), Inches(0.78)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        tf = circle.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = tag
        r.font.name = BODY_FONT
        r.font.size = Pt(13)
        r.font.bold = True
        r.font.color.rgb = PANEL

        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, Inches(3.0), Inches(1.9), Inches(1.55)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = PANEL
        card.line.color.rgb = color
        add_textbox(slide, left + Inches(0.12), Inches(3.3), Inches(1.66), Inches(0.9), label, 12.2, True, INK, BODY_FONT, PP_ALIGN.CENTER)
        connect_down(slide, circle, card, color)

        if idx < len(roadmap_items) - 1:
            connector = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                left + Inches(1.48),
                Inches(2.28),
                left + Inches(2.35),
                Inches(2.28),
            )
            connector.line.color.rgb = LINE
            connector.line.width = Pt(1.5)
            connector.line.end_arrowhead = True

    add_panel(slide, Inches(0.95), Inches(5.2), Inches(11.25), Inches(0.85), "推进原则")
    add_textbox(
        slide,
        Inches(1.15),
        Inches(5.47),
        Inches(10.8),
        Inches(0.3),
        "先让桌面本地执行闭环变深、变稳、变得可恢复，再考虑平台扩展和更多运行形态。",
        14,
        True,
        INK_SOFT,
    )
    add_footer(slide, "Recommendation: prioritize M2-quality depth over scope expansion")

    # Slide 10: Appendix
    slide = prs.slides.add_slide(title_layout)
    set_slide_background(slide)
    add_top_bar(slide, "附录：关键目录与文档入口", "需要继续跟进的人，建议从这几处开始")

    add_box_with_title(
        slide,
        Inches(0.75),
        Inches(1.55),
        Inches(4.1),
        Inches(4.75),
        "关键代码目录",
        [
            "apps/desktop/src/main",
            "apps/desktop/src/preload",
            "apps/desktop/src/renderer",
            "packages/core",
            "packages/adapters",
            "packages/contracts",
            "apps/desktop/src/main/persistence",
        ],
        accent=MANGO_DARK,
    )
    add_box_with_title(
        slide,
        Inches(5.08),
        Inches(1.55),
        Inches(3.45),
        Inches(4.75),
        "关键文档",
        [
            "docs/product/prd.md",
            "docs/product/v1-backlog.md",
            "docs/engineering/technical-architecture.md",
            "docs/engineering/detailed-architecture.md",
            "docs/engineering/module-contracts.md",
            "docs/engineering/storage-and-security.md",
        ],
        accent=INFO,
    )
    add_box_with_title(
        slide,
        Inches(8.78),
        Inches(1.55),
        Inches(3.45),
        Inches(4.75),
        "建议会后动作",
        [
            "确认下一阶段唯一主问题",
            "锁定 CLI 事件补强范围",
            "拆权限控制增强方案",
            "定义恢复 / 重试的验收标准",
            "安排一次架构评审复盘",
        ],
        accent=SUCCESS,
    )
    add_footer(slide, "Generated locally with python-pptx for editable PowerPoint delivery")

    return prs


def main():
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = build_presentation()
    prs.save(OUTPUT_PATH)
    print(OUTPUT_PATH)


if __name__ == "__main__":
    main()
